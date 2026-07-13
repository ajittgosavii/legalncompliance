"""
Build the PG&E Regulatory Compliance AI leadership deck.

    python docs/make_deck.py

Same discipline as the code: every number is either verified and cited, or explicitly labelled as a
MODEL built on assumptions that must be validated. There is no third category. A leadership deck is
exactly where an unlabelled assumption gets quoted back six months later as a commitment.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "PGE_Regulatory_Compliance_AI_Leadership_Deck.pptx")

INK      = RGBColor(0x0F, 0x17, 0x2A)
SLATE    = RGBColor(0x47, 0x55, 0x69)
MUTED    = RGBColor(0x64, 0x74, 0x8B)
RED      = RGBColor(0xB9, 0x1C, 0x1C)
RED_BG   = RGBColor(0xFE, 0xF2, 0xF2)
GREEN    = RGBColor(0x16, 0x65, 0x34)
GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
AMBER    = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG = RGBColor(0xFF, 0xFB, 0xEB)
ORANGE   = RGBColor(0xC2, 0x41, 0x0C)
ORANGE_BG= RGBColor(0xFF, 0xF7, 0xED)
BLUE     = RGBColor(0x1E, 0x40, 0xAF)
BLUE_BG  = RGBColor(0xEF, 0xF6, 0xFF)
GREY_BG  = RGBColor(0xF8, 0xFA, 0xFC)
LINE     = RGBColor(0xE2, 0xE8, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h, fill=None, line=None, lw=1.25, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (t, size, bold, colour) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(4)
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = colour; r.font.name = "Segoe UI"
    return tb


def title(s, main, sub=None, kicker=None):
    y = 0.42
    if kicker:
        text(s, 0.6, y, 12.2, 0.28, [(kicker.upper(), 10.5, True, MUTED)]); y += 0.32
    text(s, 0.6, y, 12.2, 0.62, [(main, 27, True, INK)])
    if sub:
        text(s, 0.6, y + 0.66, 12.2, 0.4, [(sub, 13.5, False, SLATE)])
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y + 1.12), Inches(12.13), Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE
    ln.line.fill.background(); ln.shadow.inherit = False


def card(s, x, y, w, h, head, body, fill=GREY_BG, line=LINE, headcol=INK, hs=13, bs=10.5):
    box(s, x, y, w, h, fill, line)
    text(s, x + 0.24, y + 0.16, w - 0.48, 0.34, [(head, hs, True, headcol)])
    text(s, x + 0.24, y + 0.62, w - 0.48, h - 0.78,
         [(l, bs, False, SLATE) for l in body], spacing=1.16)


def footer(s, note):
    text(s, 0.6, 6.95, 12.2, 0.3, [(note, 9.5, False, MUTED)])


def table(s, x, y, w, cols, rows, colw=None, fs=10.5, rh=0.34):
    n = len(rows) + 1
    tbl = s.shapes.add_table(n, len(cols), Inches(x), Inches(y), Inches(w), Inches(rh * n)).table
    if colw:
        for i, cw in enumerate(colw):
            tbl.columns[i].width = Inches(cw)
    for j, c in enumerate(cols):
        cell = tbl.cell(0, j); cell.text = ""
        r = cell.text_frame.paragraphs[0].add_run(); r.text = c
        r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
        cell.fill.solid(); cell.fill.fore_color.rgb = INK
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = ""
            r = cell.text_frame.paragraphs[0].add_run(); r.text = str(v)
            r.font.size = Pt(fs); r.font.color.rgb = SLATE; r.font.name = "Segoe UI"
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else GREY_BG
    return tbl


# ============================== 1 TITLE ==============================
s = slide()
box(s, 0, 0, 13.333, 7.5, INK, None, radius=False)
text(s, 0.9, 1.85, 11.5, 0.6, [("Regulatory Compliance AI", 46, True, WHITE)])
text(s, 0.9, 2.85, 11.5, 0.45,
     [("From reacting to regulatory change — to getting ahead of it.", 21, False, RGBColor(0xCB,0xD5,0xE1))])
box(s, 0.9, 3.65, 5.4, 0.05, ORANGE, None, radius=False)
text(s, 0.9, 4.0, 11.5, 1.3, [
    ("Seven agents turn dense regulatory prose into atomic, testable obligations — each with a", 15, False, RGBColor(0x94,0xA3,0xB8)),
    ("machine-verified citation back to the source — score their impact, and assemble audit-ready", 15, False, RGBColor(0x94,0xA3,0xB8)),
    ("evidence packages with gap analysis. A human always signs off.", 15, False, RGBColor(0x94,0xA3,0xB8)),
], spacing=1.25)
text(s, 0.9, 6.45, 11.5, 0.4,
     [("Leadership briefing  ·  Infosys  ·  Anchor client: Pacific Gas and Electric Company", 12, False, MUTED)])

# ============================== 2 THE PROBLEM ==============================
s = slide()
title(s, "The problem", "PG&E answers to nine regulators. Today, proving compliance is manual.",
      kicker="Problem")

card(s, 0.6, 1.9, 5.95, 2.3, "What happens today",
     ["• Analysts read hundreds of pages of dense legal prose",
      "• Obligations are hand-extracted into spreadsheets",
      "• Document owners are chased for audit evidence",
      "• Enforcement precedent lives in institutional memory",
      "• Evidence gaps are discovered DURING an audit — not before it"],
     GREY_BG, LINE, INK)

card(s, 6.78, 1.9, 5.95, 2.3, "Why PG&E specifically",
     ["PG&E's binding constraint is not KNOWING the rule.",
      "It is PROVING compliance on demand, repeatedly, under sustained regulatory scrutiny.",
      "So the costly failure mode is not a rule discovered late. It is an evidence gap found by an auditor rather than by PG&E."],
     BLUE_BG, RGBColor(0xBF,0xDB,0xFE), BLUE)

box(s, 0.6, 4.45, 12.13, 1.5, ORANGE_BG, ORANGE, 2)
text(s, 0.9, 4.62, 11.6, 1.25, [
    ("The insight that makes this a business case, not a productivity tool", 15, True, ORANGE),
    ("Under AB 1054, demonstrated WMP compliance and the annual Safety Certification affect the presumption applied in", 12, False, SLATE),
    ("cost-recovery proceedings. Separately, the CPUC can DISALLOW claimed costs where compliance cannot be evidenced.", 12, False, SLATE),
    ("Evidence quality is therefore an INPUT TO COST RECOVERY — not administrative overhead.", 13, True, INK),
], spacing=1.18)

footer(s, "⚠ The precise legal mechanics of the AB 1054 presumption must be confirmed with qualified counsel before any dollar figure is attached to this chain. What is defensible is the mechanism — not yet a quantified saving.")

# ============================== 3 WHAT IT IS ==============================
s = slide()
title(s, "What the platform is", "Four modules. Seven agents. One traceable chain.",
      kicker="Description")

mods = [
    ("Module 1 — Regulatory Monitor", "Agentic · 5-node pipeline",
     "Monitors 9 regulators. Classifies each change, then extracts obligations: who must do what, by when, measured how, penalty if not — each with a VERIFIED source quote."),
    ("Module 2 — Obligation Impact", "Agentic · 4-node graph",
     "Decomposes a regulation into atomic obligations, finds conflicts with the existing estate, and scores four dimensions: cost, operations, timeline, penalty risk."),
    ("Module 3 — Audit Preparation", "Agentic · Supervisor + 3 specialists",
     "THE FLAGSHIP. Evidence Collector, Gap Analyzer and Response Drafter, coordinated by a Supervisor. Ends in a 0–100 audit readiness score."),
    ("Module 4 — Case Analytics", "Generative · RAG",
     "Precedent search and penalty-trend analysis over the enforcement record. Loads with NO LLM call — the zero-cost entry point for any demo."),
]
for i, (h, t, b) in enumerate(mods):
    x = 0.6 + (i % 2) * 6.22
    y = 1.95 + (i // 2) * 1.75
    fill = ORANGE_BG if "Module 3" in h else GREY_BG
    ln = ORANGE if "Module 3" in h else LINE
    hc = ORANGE if "Module 3" in h else INK
    box(s, x, y, 5.98, 1.6, fill, ln, 1.75 if "Module 3" in h else 1.25)
    text(s, x + 0.24, y + 0.14, 5.5, 0.3, [(h, 13, True, hc)])
    text(s, x + 0.24, y + 0.44, 5.5, 0.24, [(t, 9.5, True, MUTED)])
    text(s, x + 0.24, y + 0.72, 5.5, 0.8, [(b, 10.5, False, SLATE)], spacing=1.16)

box(s, 0.6, 5.55, 12.13, 1.15, INK, None)
text(s, 0.9, 5.72, 11.6, 0.9, [
    ("What Module 3 actually produces — the artefact a regulator asks for when it says “show me”:", 12, False, RGBColor(0xCB,0xD5,0xE1)),
    ("REGULATION → OBLIGATION → EVIDENCE DOCUMENT → GAP → OWNER + DATE → DRAFTED RESPONSE → READINESS SCORE", 13.5, True, WHITE),
], spacing=1.2)

# ============================== 4 REGULATORS + THE OEIS FIX ==============================
s = slide()
title(s, "The regulatory perimeter — and the detail that earns credibility",
      "Nine bodies. One of them was missing, and its absence would have ended the conversation.",
      kicker="Domain grounding")

regs = ["CPUC", "OEIS", "FERC", "NERC", "CARB", "EPA", "PHMSA", "Cal-OSHA", "CEC"]
for i, rr in enumerate(regs):
    x = 0.6 + i * 1.36
    is_oeis = rr == "OEIS"
    box(s, x, 1.95, 1.2, 0.55, ORANGE_BG if is_oeis else BLUE_BG,
        ORANGE if is_oeis else RGBColor(0xBF,0xDB,0xFE), 2 if is_oeis else 1.25)
    text(s, x + 0.05, 2.08, 1.1, 0.3, [(rr, 11.5, True, ORANGE if is_oeis else BLUE)],
         align=PP_ALIGN.CENTER)

box(s, 0.6, 2.75, 12.13, 1.75, ORANGE_BG, ORANGE, 2)
text(s, 0.9, 2.92, 11.6, 1.5, [
    ("A defect we found and fixed — and it matters more than it looks", 15, True, ORANGE),
    ("The original codebase routed Wildfire Mitigation Plans to the CPUC, and did not list OEIS as a regulator at all.", 12, False, SLATE),
    ("Since 2021, WMPs are filed with and approved by the Office of Energy Infrastructure Safety (OEIS), which also issues", 12, False, SLATE),
    ("the annual Safety Certification. The CPUC ratifies.", 12, False, SLATE),
    ("Getting this wrong in front of PG&E's regulatory affairs team — on the single topic they care most about — would have", 12, True, INK),
    ("ended the credibility conversation before the demo started.", 12, True, INK),
], spacing=1.16)

card(s, 0.6, 4.75, 5.95, 1.75, "What we changed",
     ["• OEIS is now a first-class monitored regulator",
      "• The WMP regulation is retargeted to OEIS",
      "• The system prompt carries an explicit instruction: “Never state that a WMP is filed with or approved by the CPUC.”"],
     GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN)

card(s, 6.78, 4.75, 5.95, 1.75, "Why it is the whole argument in miniature",
     ["This is the error a GENERIC RegTech tool makes and a DOMAIN-GROUNDED one does not.",
      "Regulatory nuance is not a feature. It is the difference between a tool PG&E trusts and one it politely declines."],
     GREY_BG, LINE, INK)

# ============================== 5 PAIN AREAS ==============================
s = slide()
title(s, "The pain this removes", "Where the hours go, and where the money leaks",
      kicker="Pain areas")

pains = [
    ("Hours of manual rule-reading", "A single CPUC decision contains a dozen separately-testable duties buried in prose, each with its own deadline, measurement method and penalty. An analyst extracts them by hand into a spreadsheet."),
    ("Audit-day surprises", "Evidence gaps are found when the auditor finds them. By then the only options are bad ones."),
    ("Stale evidence, invisibly", "A procedure that was current three years ago still looks current in the folder. Nobody re-reads 48 documents before an audit."),
    ("Conflicts nobody catches", "New OEIS and CPUC requirements routinely collide with GO 95/165 and existing WMP commitments. Nothing systematically checks."),
    ("Precedent as folklore", "Penalty-exposure conversations rest on who remembers what. Decades of enforcement history sit unread."),
    ("Deadlines discovered late", "An obligation with no owner and no date is a missed deadline waiting for a calendar."),
]
for i, (h, b) in enumerate(pains):
    x = 0.6 + (i % 2) * 6.22
    y = 1.95 + (i // 2) * 1.62
    card(s, x, y, 5.98, 1.45, h, [b], GREY_BG, LINE, INK, hs=12.5, bs=10.5)

# ============================== 6 END-TO-END FLOW ==============================
s = slide()
title(s, "End-to-end flow", "From regulator to signed-off audit package", kicker="How it works")

steps = [
    ("1", "INGEST", "Monitor 9 regulators\nfor change", BLUE_BG, BLUE),
    ("2", "EXTRACT", "Atomic obligations +\nVERBATIM source quote", ORANGE_BG, ORANGE),
    ("3", "VERIFY", "Does that quote really\nexist? Checked in CODE", RED_BG, RED),
    ("4", "SCORE", "Cost · operations ·\ntimeline · penalty risk", AMBER_BG, AMBER),
    ("5", "ASSEMBLE", "Evidence → gaps →\nowners → responses", GREEN_BG, GREEN),
    ("6", "SIGN OFF", "Human reviewer.\nThe platform files nothing", INK, WHITE),
]
for i, (n, h, b, fill, col) in enumerate(steps):
    x = 0.6 + i * 2.05
    box(s, x, 2.1, 1.85, 2.2, fill, LINE if fill != INK else INK)
    text(s, x + 0.15, 2.22, 1.55, 0.3, [(n, 15, True, col)])
    text(s, x + 0.15, 2.6, 1.55, 0.3, [(h, 12, True, col)])
    text(s, x + 0.15, 3.05, 1.55, 1.1,
         [(l, 9.5, False, col if fill == INK else SLATE) for l in b.split("\n")], spacing=1.12)
    if i < 5:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.87), Inches(3.0),
                                Inches(0.16), Inches(0.22))
        ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(0xCB,0xD5,0xE1)
        ar.line.fill.background(); ar.shadow.inherit = False

box(s, 0.6, 4.65, 12.13, 1.0, RED_BG, RED, 1.75)
text(s, 0.85, 4.82, 11.6, 0.75, [
    ("THE GATE — provenance is a CODE check, not a prompt instruction", 13, True, RED),
    ("Every obligation must quote the regulation verbatim, and the system verifies the quote exists. The model cannot talk its way past code. An obligation it cannot trace is flagged ⚠ UNVERIFIED — never presented as fact.", 11, False, SLATE),
], spacing=1.15)

box(s, 0.6, 5.9, 12.13, 0.85, GREEN_BG, RGBColor(0x86,0xEF,0xAC), 1.75)
text(s, 0.85, 6.05, 11.6, 0.6, [
    ("Every package is stamped PENDING_HUMAN_REVIEW. The platform PREPARES the submission. It never files it. PG&E remains accountable.", 12.5, True, GREEN),
])

# ============================== 7 ARCHITECTURE ==============================
s = slide()
title(s, "High-level architecture", "The agents contain no industry logic — that is the reuse thesis.",
      kicker="Architecture")

box(s, 0.6, 1.9, 12.13, 0.75, BLUE_BG, RGBColor(0xBF,0xDB,0xFE), 1.4)
text(s, 0.8, 1.98, 3.5, 0.24, [("1 — REGULATORY SOURCES · 9 BODIES", 9.5, True, MUTED)])
regs = ["CPUC", "OEIS", "FERC", "NERC", "CARB", "EPA", "PHMSA", "Cal-OSHA", "CEC"]
for i, rr in enumerate(regs):
    x = 0.85 + i * 1.31
    is_o = rr == "OEIS"
    box(s, x, 2.24, 1.2, 0.32, ORANGE_BG if is_o else WHITE,
        ORANGE if is_o else RGBColor(0xBF,0xDB,0xFE), 1.5 if is_o else 1)
    text(s, x + 0.03, 2.29, 1.14, 0.24, [(rr, 9, True, ORANGE if is_o else BLUE)], align=PP_ALIGN.CENTER)

box(s, 0.6, 2.8, 12.13, 1.6, AMBER_BG, RGBColor(0xFC,0xD3,0x4D), 1.4)
text(s, 0.8, 2.88, 6.0, 0.24, [("2 — SEVEN AGENTS · FOUR WORKFLOWS (LangGraph)", 9.5, True, MUTED)])
wfs = [
    (0.85, 2.72, "WF-01 Monitor", "Fetch → Classify → Extract+VERIFY → Map → Alert", WHITE, LINE, INK),
    (3.72, 2.35, "WF-02 Impact", "Decompose → Cross-ref → Score → Report", WHITE, LINE, INK),
    (6.22, 3.6, "WF-03 Audit Prep — FLAGSHIP", "Supervisor + Evidence + Gaps + Responses → 0–100 score", ORANGE_BG, ORANGE, ORANGE),
    (10.0, 2.6, "WF-04 Cases", "RAG precedent · trends · risk", WHITE, LINE, INK),
]
for (x, w, h, b, f, l, c) in wfs:
    box(s, x, 3.16, w, 1.05, f, l, 1.75 if c == ORANGE else 1.2)
    text(s, x + 0.12, 3.24, w - 0.24, 0.26, [(h, 10.5, True, c)])
    text(s, x + 0.12, 3.52, w - 0.24, 0.6, [(b, 9, False, SLATE)], spacing=1.1)

box(s, 0.6, 4.55, 5.95, 1.15, RED_BG, RED, 1.75)
text(s, 0.8, 4.65, 5.6, 0.95, [
    ("3a — PROVENANCE: a CODE check, not a prompt", 11, True, RED),
    ("Every obligation must quote the regulation verbatim; the system verifies the quote exists in the source. Unverifiable → flagged, downgraded, never shown as fact.", 9.5, False, SLATE),
], spacing=1.1)

box(s, 6.78, 4.55, 5.95, 1.15, GREEN_BG, RGBColor(0x86,0xEF,0xAC), 1.75)
text(s, 6.98, 4.65, 5.6, 0.95, [
    ("3b — HUMAN REVIEW GATE", 11, True, GREEN),
    ("PENDING_HUMAN_REVIEW on every package. Decision support, never a compliance determination. Nothing reaches a regulator without sign-off.", 9.5, False, SLATE),
], spacing=1.1)

box(s, 0.6, 5.85, 12.13, 0.85, INK, None)
text(s, 0.85, 5.98, 11.6, 0.65, [
    ("4 — SWAPPABLE INDUSTRY PACKS   ·   Energy & Utilities (PG&E) · Retail · Resources · Services", 12, True, WHITE),
    ("Switch the pack and all 7 agents re-target. No agent code changes. Each new vertical ≈ 2.5 FTE-months.", 10, False, RGBColor(0xCB,0xD5,0xE1)),
], spacing=1.15)

footer(s, "Editable SVG delivered separately: docs/RegCompliance_Architecture.svg")

# ============================== 8 ADVANTAGES / STANDS OUT ==============================
s = slide()
title(s, "How this stands out", "We went looking for a reason not to build it. This is what survived.",
      kicker="Differentiation")

box(s, 0.6, 1.9, 5.95, 1.75, GREY_BG, LINE, 1.4)
text(s, 0.85, 2.02, 5.5, 1.5, [
    ("Be ready for this question", 12.5, True, INK),
    ("“Doesn't Archer / Workiva / MetricStream already do obligation extraction?”", 11, False, SLATE),
    ("Yes — for a GENERIC, cross-industry compliance market.", 11, True, INK),
], spacing=1.16)

box(s, 6.78, 1.9, 5.95, 1.75, GREEN_BG, RGBColor(0x86,0xEF,0xAC), 1.75)
text(s, 7.03, 2.02, 5.5, 1.5, [
    ("The answer", 12.5, True, GREEN),
    ("This is not a product competing for market share. It is a CLIENT IMPLEMENTATION in PG&E's environment, grounded in PG&E's regulators, PG&E's departments, and PG&E's obligation estate.", 11, False, SLATE),
    ("A generic tool routes WMPs to the CPUC. We do not.", 11, True, INK),
], spacing=1.16)

diffs = [
    ("Provenance is CODE, not a prompt", "The system verifies the model's citation against the real regulation. Unverifiable obligations are flagged, never presented as fact. This is what makes it deployable in a regulated utility at all."),
    ("Human review is a gate", "Module 3 drafts submission-grade text. Every package is stamped PENDING_HUMAN_REVIEW. The platform prepares; it never files."),
    ("Deep domain grounding", "OEIS vs CPUC. GO 95. GO 165. Rule 20. HFTD tiers. PSPS. This is the vocabulary PG&E actually uses — and getting it wrong is disqualifying."),
    ("Swappable industry packs", "The same 7 agents serve Retail, Resources and Services. Build once, re-target. Each new vertical ≈ 2.5 FTE-months."),
]
for i, (h, b) in enumerate(diffs):
    x = 0.6 + (i % 2) * 6.22
    y = 3.85 + (i // 2) * 1.5
    card(s, x, y, 5.98, 1.35, h, [b], GREY_BG, LINE, INK, hs=12, bs=10)

footer(s, "Honest limit: features are not a moat. The durable advantages are domain grounding, the provenance method, integration depth into PG&E's document systems, and the client relationship.")

# ============================== 9 FINANCIAL ==============================
s = slide()
title(s, "Financial model — Infosys",
      "Every figure here is ILLUSTRATIVE and must be validated in discovery. None of it is a finding.",
      kicker="Commercials")

box(s, 0.6, 1.85, 12.13, 0.62, AMBER_BG, RGBColor(0xFC,0xD3,0x4D), 1.5)
text(s, 0.85, 1.95, 11.6, 0.45, [
    ("READ THIS FIRST — these are ASSUMPTION PLACEHOLDERS, not measurements. A leadership deck is exactly where an unlabelled", 10.5, True, AMBER),
    ("assumption gets quoted back six months later as a commitment. What must be validated is named at the bottom of this slide.", 10.5, False, SLATE),
], spacing=1.1)

table(s, 0.6, 2.65, 7.3,
      ["Stage", "Duration", "Indicative revenue"],
      [["Qualify & discovery", "2–4 weeks", "$100k – 200k"],
       ["Single-audit pilot", "6–8 weeks", "$300k – 600k"],
       ["Production build (ingestion, provenance, workflow)", "6–9 months", "$1.2M – 2.5M"],
       ["Managed run", "annual", "$300k – 800k ARR"]],
      colw=[4.2, 1.5, 1.6], fs=10.5)

card(s, 8.2, 2.65, 4.53, 1.15, "Anchor engagement (PG&E)",
     ["~$2M – 4M over 2–3 years"], BLUE_BG, RGBColor(0xBF,0xDB,0xFE), BLUE, hs=11.5, bs=13)

card(s, 8.2, 3.95, 4.53, 1.35, "Run cost is NOT the constraint",
     ["A full 11-regulation monitoring run: well under $1. An audit-prep run: ~$0.20.",
      "The cost is engineering, integration and change management."],
     GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN, hs=11.5, bs=10)

card(s, 0.6, 4.55, 7.3, 1.25, "Where the value actually is — and the order to pitch it",
     ["1. AUDIT READINESS (lead here). 2. Obligation decomposition. 3. Conflict detection. 4. Precedent.",
      "Do NOT lead with regulatory monitoring — it is the least differentiated capability. PG&E's regulatory affairs function already tracks CPUC dockets closely."],
     GREY_BG, LINE, INK, hs=12, bs=10)

card(s, 8.2, 5.45, 4.53, 0.95, "MUST be validated first",
     ["• FTE-hours spent on obligation extraction today",
      "• Cost of responding to ONE audit / data request"],
     AMBER_BG, RGBColor(0xFC,0xD3,0x4D), AMBER, hs=11, bs=9.5)

footer(s, "We deliberately do NOT anchor value to PG&E's headline penalty history. Those penalties came from OPERATIONAL failures — pipeline integrity, vegetation management. This platform does not prevent a wildfire. Claiming a share of them would be dishonest, and a PG&E audience would see through it instantly.")

# ============================== 10 EFFORT ==============================
s = slide()
title(s, "Implementation effort", "What exists today, and what production actually costs",
      kicker="Effort")

card(s, 0.6, 1.9, 5.95, 2.15, "Already built and verified",
     ["• 7 agents, 4 workflows — running end-to-end",
      "• Provenance verification wired into obligation extraction",
      "• Human-review gate (PENDING_HUMAN_REVIEW)",
      "• 4 industry packs — all build 7/7 agent prompts",
      "• 9 regulators incl. OEIS; 11 regs · 29 cases · 48 evidence docs",
      "• All 5 pages render across every pack"],
     GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN)

table(s, 6.78, 1.9, 5.95,
      ["Gap to close", "Effort"],
      [["Live regulator ingestion (CPUC/OEIS/FERC)", "L — 8–16 wks"],
       ["Real embeddings + vector store", "S — 1–3 wks"],
       ["Persistence + obligation lifecycle", "M — 3–8 wks"],
       ["Span-level provenance (offset into the PDF)", "M"],
       ["Evaluation harness (golden set)", "M"],
       ["SSO / RBAC / access-scoped retrieval", "M"]],
      colw=[4.35, 1.6], fs=10)

box(s, 0.6, 4.35, 12.13, 0.95, INK, None)
text(s, 0.85, 4.52, 11.6, 0.7, [
    ("Production MVP: ~7–9 months elapsed  ·  ~28 FTE-months  ·  indicative build $1.2M – 2.5M", 14, True, WHITE),
    ("Each additional industry pack thereafter: ~2.5 FTE-months. That reuse ratio is the strongest commercial argument here.", 11, False, RGBColor(0xCB,0xD5,0xE1)),
], spacing=1.2)

card(s, 0.6, 5.55, 12.13, 1.2, "The biggest single work item — and it is not the AI",
     ["LIVE INGESTION. Today the regulations are fixture data; ingestion/scrapers/ and ingestion/parsers/ are empty stub packages. Production needs real connectors to CPUC, OEIS and FERC publication feeds, PDF and HTML parsing, change detection and de-duplication.",
      "Do not present this as a system that monitors anything today. It is a working prototype of the target architecture."],
     RED_BG, RED, RED, hs=12.5, bs=10)

# ============================== 11 WHAT IS PROVEN ==============================
s = slide()
title(s, "What is proven — and what is not",
      "Stated before anyone has to ask. This is the slide that makes the other ten believable.",
      kicker="Credibility")

card(s, 0.6, 1.9, 5.95, 2.4, "Verified by test",
     ["• All 4 industry packs build 7/7 agent prompts",
      "• All 5 pages render across every pack",
      "• Provenance verifier is wired and reports citations verified per run",
      "• The Supervisor's plan now reaches all 3 specialists (a real bug we found and fixed)",
      "• KPI counts computed from the data — they previously claimed 12/28/$19.8B against real corpora of 11/29/$17.47B"],
     GREEN_BG, RGBColor(0x86,0xEF,0xAC), GREEN, hs=13, bs=10)

card(s, 6.78, 1.9, 5.95, 2.4, "NOT yet measured",
     ["• AGENT OUTPUT QUALITY IS UNMEASURED. There is no golden set and no evaluation harness. Until one exists, every claim about extraction accuracy is an anecdote.",
      "• How often GPT-4o quotes VERBATIM rather than paraphrasing is unknown. If it paraphrases, obligations correctly show as UNVERIFIED — the system stays honest, but the review looks weak."],
     AMBER_BG, RGBColor(0xFC,0xD3,0x4D), AMBER, hs=13, bs=10)

box(s, 0.6, 4.5, 12.13, 2.2, RED_BG, RED, 1.75)
text(s, 0.85, 4.65, 11.6, 1.95, [
    ("What we will NOT claim — we say this first, not when challenged", 14, True, RED),
    ("•  NO LIVE INGESTION. The regulations are fixture data. “Continuously monitors 9 regulators” is the TARGET architecture, not today's function.", 11, False, SLATE),
    ("•  All corpora are ILLUSTRATIVE. Before a PG&E conversation the case corpus must be replaced with the real public enforcement record — PG&E will recognise its own history instantly.", 11, False, SLATE),
    ("•  THIS PLATFORM DOES NOT PREVENT A WILDFIRE. PG&E's headline penalties came from operational failures. We do not claim a share of them.", 11, True, INK),
    ("•  Obligations carry no span-level citation into the source PDF yet. For a compliance register that is required, and it is on the roadmap.", 11, False, SLATE),
], spacing=1.16)

footer(s, "Owning the boundary makes every other claim more believable — and it is the only position that survives Q&A with people who know this domain.")

# ============================== 12 THE ASK ==============================
s = slide()
box(s, 0, 0, 13.333, 7.5, INK, None, radius=False)
text(s, 0.9, 0.9, 11.5, 0.5, [("The recommendation", 30, True, WHITE)])
box(s, 0.9, 1.6, 3.4, 0.05, ORANGE, None, radius=False)
text(s, 0.9, 1.95, 11.5, 0.5,
     [("Do not pitch a platform. Pitch one narrow, falsifiable proof.", 19, True, RGBColor(0xCB,0xD5,0xE1))])

steps = [
    ("QUALIFY", "weeks 0–2", "Discovery with Regulatory Affairs, Internal Audit and Compliance. Size the value with PG&E's OWN numbers. Pick ONE upcoming audit or data request as the target."),
    ("PROVE", "weeks 2–8", "A single-audit pilot. Load the REAL obligations and the REAL evidence inventory for that audit's scope. Run Module 3."),
    ("EXPAND", "post-pilot", "Only if the pilot clears the bar: live ingestion for CPUC + OEIS, span-level provenance, the review workflow, persistence."),
]
for i, (h, w, b) in enumerate(steps):
    y = 2.95 + i * 1.2
    box(s, 0.9, y, 11.5, 1.05, RGBColor(0x1E, 0x29, 0x3B), None)
    text(s, 1.15, y + 0.12, 2.0, 0.3, [(h, 13, True, WHITE)])
    text(s, 1.15, y + 0.47, 2.0, 0.28, [(w, 10, False, RGBColor(0x94,0xA3,0xB8))])
    text(s, 3.3, y + 0.18, 8.9, 0.8, [(b, 11.5, False, RGBColor(0xCB,0xD5,0xE1))], spacing=1.15)

box(s, 0.9, 6.6, 11.5, 0.6, RGBColor(0x45, 0x14, 0x14), None)
text(s, 1.15, 6.72, 11.0, 0.4, [
    ("Success is BINARY: did Module 3 surface a real evidence gap the team did not already know about? The pilot is designed so that it CAN FAIL — and if it does, we say so.", 11.5, True, RGBColor(0xFC,0xA5,0xA5)),
])

prs.save(OUT)
print("WROTE:", OUT)
